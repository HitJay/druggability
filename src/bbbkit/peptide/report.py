"""
bbbkit.peptide.report — BBB 预测自动报告（ESM-2 预测 → Opus 解读 → 自包含 HTML）

把"预测结果 → 适用域标记 → LLM 叙述 → 渲染 HTML"串成一步，面向肽类 BBB 通透性。

复用 ``bbbkit.report.llm.LLMClient``（OpenAI 兼容，已对接 .env 的 MARKETPLACE_*，
默认 anthropic_claude_opus_4_7）。新增的是 BBB/肽专用层：

- ``applicability_domain``  适用域标记（长度 / 修饰 / MW 规则；基于 B3Pred 实测分布）
- ``confidence_from_ci``    由 bootstrap 置信区间宽度给可信度等级
- ``generate_peptide_narrative`` / ``generate_exec_summary``  Opus 解读（优雅降级到模板）
- ``render_html``          自包含 HTML（navy 主题 + CSS 条形图 + 诚实口径）
- ``build_bbb_report``     端到端编排

设计原则（沿用 llm.py 的"优雅降级"）：无 key / 无网络 / SDK 缺失时，所有 LLM 调用
回退到基于规则的模板，报告始终能产出。适用域与 HTML 渲染是纯 Python（零第三方依赖）。

诚实口径贯穿全报告：模型预测的是 *内在、基于序列的通透倾向（propensity）*，
而非体内脑暴露或转运路径；out-of-domain 输入（>20–30 aa 或脂化/PEG/环化/D-aa）
的分数只是"骨架上界"，不是可信判定。
"""

from __future__ import annotations

import datetime as _dt
import html as _html
import json
import logging
import math
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ── 适用域常量（来自 B3Pred Dataset-1 训练集实测：405 肽、中位 13 aa、80% ≤20、max 30）──
LEN_IN_DOMAIN = 20      # ≤ 此长度视为域内
LEN_EXTRAPOLATION = 30  # 训练集最大长度；> 此为明确外推
MW_LARGE_WARN = 1500.0  # 大分子告警阈（描述符肽工具不可靠区）

# ESM-2 无法编码的修饰关键词（命中即 out-of-domain-modification）
_MOD_PATTERNS = re.compile(
    r"lipid|acyl|palmit|stear/?|c1[68]|c20|peg|pegylat|diacid|"
    r"cycli|disulf|staple|lactam|d-?(ala|phe|amino)|n-?methyl|amidat",
    re.IGNORECASE,
)

# 非肽 / 小分子标识（→ 走小分子专用工具 B3clf，而非 ESM 序列模型）
_MODALITY_PATTERNS = re.compile(
    r"small[\s-]?molecule|non[\s-]?peptide|smiles|\bsm\b|b3clf|cns-?mpo",
    re.IGNORECASE,
)
# 序列模型确实无能为力、且无专用工具的模态（标 OOD-modality）
_UNSUPPORTED_MODALITY = re.compile(
    r"oligonucleotide|antibody|\bnce\b|\bsiRNA\b|\bmRNA\b",
    re.IGNORECASE,
)

# B3clf 小分子工具 benchmark（既往验证：12-model 共识 70%，XGBoost ADASYN 72.7%）
B3CLF_BENCHMARK = "B3clf small-molecule benchmark ~70-73% (12-model consensus 70%, XGBoost-ADASYN 72.7%)"


def applicability_domain(
    length: int | None,
    modification: str | None = None,
    mw: float | None = None,
    method: str | None = None,
) -> tuple[str, str]:
    """返回 (domain, reason)。

    domain ∈ {in-domain, edge-extrapolation, out-of-domain-length,
              out-of-domain-modification, out-of-domain-modality, small-molecule}。
    优先级：method/modality 判定 > modification（修饰）> length（长度）。

    小分子（method=b3clf/small_molecule 或 modification 含小分子标识）→ ``small-molecule``：
    这是一个 **有效** 类别，预测来自专用工具 B3clf（已 benchmark），而非 ESM 序列模型；
    不应被当作 out-of-domain/不可信。
    """
    reasons: list[str] = []

    meth = (method or "").strip().lower()
    mod = (modification or "").strip()
    is_sm = meth in {"b3clf", "small_molecule", "small-molecule", "sm", "cns-mpo"}
    if not is_sm and mod and mod.lower() not in {"none", "native", "-", "l-backbone"}:
        is_sm = bool(_MODALITY_PATTERNS.search(mod))

    if is_sm:
        return (
            "small-molecule",
            f"Non-peptide small molecule: predicted by the dedicated small-molecule "
            f"tool (B3clf / CNS-MPO), not the ESM-2 sequence model. {B3CLF_BENCHMARK}.",
        )

    if mod and mod.lower() not in {"none", "native", "-", "l-backbone"}:
        # 无专用工具的非肽模态（抗体/寡核苷酸…）→ 真正 OOD
        if _UNSUPPORTED_MODALITY.search(mod):
            return (
                "out-of-domain-modality",
                f"Input '{mod}' is neither a standard peptide nor covered by the "
                f"small-molecule tool; no validated BBB predictor applies here.",
            )
        if _MOD_PATTERNS.search(mod):
            return (
                "out-of-domain-modification",
                f"Modification '{mod}' cannot be encoded by ESM-2 (20 standard "
                f"amino acids only); score is a backbone-level upper bound.",
            )

    if mw is not None and mw > MW_LARGE_WARN:
        reasons.append(
            f"MW {mw:.0f} Da exceeds the reliable range of descriptor peptide "
            f"tools (>{MW_LARGE_WARN:.0f})"
        )

    if length is None:
        dom = "in-domain"
    elif length <= LEN_IN_DOMAIN:
        dom = "in-domain"
        reasons.insert(0, f"Length {length} aa within training distribution (B3Pred 80% <=20 aa)")
    elif length <= LEN_EXTRAPOLATION:
        dom = "edge-extrapolation"
        reasons.insert(0, f"Length {length} aa near training maximum (B3Pred max 30 aa)")
    else:
        dom = "out-of-domain-length"
        reasons.insert(0, f"Length {length} aa beyond training maximum (B3Pred max 30 aa) -> extrapolation")

    # MW 告警可把 in-domain 降级为 edge
    if dom == "in-domain" and mw is not None and mw > MW_LARGE_WARN:
        dom = "edge-extrapolation"

    return dom, "; ".join(reasons) if reasons else "-"


def confidence_from_ci(ci_low: float | None, ci_high: float | None) -> tuple[str, float | None]:
    """由 90% bootstrap 区间宽度给可信度等级。返回 (level, width)。"""
    if ci_low is None or ci_high is None:
        return "unknown", None
    width = abs(ci_high - ci_low)
    if width < 30:
        level = "high"
    elif width < 60:
        level = "medium"
    else:
        level = "low"
    return level, width


def _parse_ci(ci: Any) -> tuple[float | None, float | None]:
    """把 '[5.7, 99.8]' 这类字符串解析为 (low, high)。"""
    if ci is None:
        return None, None
    if isinstance(ci, (list, tuple)) and len(ci) == 2:
        try:
            return float(ci[0]), float(ci[1])
        except (TypeError, ValueError):
            return None, None
    nums = re.findall(r"-?\d+\.?\d*", str(ci))
    if len(nums) >= 2:
        return float(nums[0]), float(nums[1])
    return None, None


def _parse_committee_votes(d: dict[str, Any], known_route: str) -> tuple[int | None, int | None]:
    """Parse committee votes from explicit fields or free text like '5/12 BBB+'."""
    for a, b in (("vote_pos", "vote_total"), ("committee_pos", "committee_n")):
        if d.get(a) not in (None, "") and d.get(b) not in (None, ""):
            try:
                pos = int(float(d[a]))
                n = int(float(d[b]))
                if 0 <= pos <= n and n > 0:
                    return pos, n
            except (TypeError, ValueError):
                pass

    packed = d.get("committee_votes", d.get("consensus_votes", ""))
    m = re.search(r"(\d+)\s*/\s*(\d+)", str(packed))
    if m:
        pos, n = int(m.group(1)), int(m.group(2))
        if 0 <= pos <= n and n > 0:
            return pos, n

    m = re.search(r"(\d+)\s*/\s*(\d+)\s*BBB\+", known_route, re.IGNORECASE)
    if m:
        pos, n = int(m.group(1)), int(m.group(2))
        if 0 <= pos <= n and n > 0:
            return pos, n
    return None, None


def _wilson_interval(pos: int, n: int, z: float = 1.645) -> tuple[float, float]:
    """Wilson score interval for binomial proportion (default 90% with z=1.645)."""
    if n <= 0:
        return 0.0, 1.0
    phat = pos / n
    denom = 1.0 + (z * z) / n
    center = (phat + (z * z) / (2.0 * n)) / denom
    half = (z / denom) * math.sqrt((phat * (1.0 - phat) / n) + ((z * z) / (4.0 * n * n)))
    lo = max(0.0, center - half)
    hi = min(1.0, center + half)
    return lo, hi


# ─────────────────────────────────────────────────────────────────────────────
# 记录规范化
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class PeptideRecord:
    """单条肽的 BBB 预测记录（报告内部规范结构）。"""

    name: str
    sequence: str = ""
    length: int | None = None
    modification: str = ""
    p_bbb: float | None = None          # P(BBB+) 百分数 0–100
    call: str = ""                       # BBB+ / BBB-
    ci_low: float | None = None
    ci_high: float | None = None
    known_route: str = ""                # 文献生物学路由（[LITERATURE]）
    receptor: str = ""
    method: str = "esm"                  # 预测来源：esm(默认肽序列) / b3clf(小分子)
    domain: str = ""
    domain_reason: str = ""
    confidence: str = ""
    ci_width: float | None = None
    committee_pos: int | None = None
    committee_n: int | None = None
    narrative: str = ""

    @classmethod
    def from_dict(cls, d: dict[str, Any]) -> "PeptideRecord":
        seq = str(d.get("sequence", "") or "")
        length = d.get("length")
        if length in (None, "") and seq:
            length = len(seq)
        try:
            length = int(length) if length not in (None, "") else None
        except (TypeError, ValueError):
            length = len(seq) or None

        p = d.get("p_bbb", d.get("ESM2_P(BBB+)%", d.get("P(BBB+)%")))
        try:
            p = float(p) if p not in (None, "") else None
        except (TypeError, ValueError):
            p = None

        ci_low, ci_high = _parse_ci(d.get("ci") or d.get("ESM2_boot_90CI"))
        if d.get("ci_low") is not None:
            ci_low = float(d["ci_low"])
        if d.get("ci_high") is not None:
            ci_high = float(d["ci_high"])

        rec = cls(
            name=str(d.get("name", d.get("id", "peptide"))),
            sequence=seq,
            length=length,
            modification=str(d.get("modification", "") or ""),
            p_bbb=p,
            call=str(d.get("call", d.get("ESM2_call", "")) or ""),
            ci_low=ci_low,
            ci_high=ci_high,
            known_route=str(d.get("known_route", d.get("known_central_route", "")) or ""),
            receptor=str(d.get("receptor", "") or ""),
            method=str(d.get("method", "") or "esm").strip() or "esm",
        )
        rec.committee_pos, rec.committee_n = _parse_committee_votes(d, rec.known_route)
        rec.domain, rec.domain_reason = applicability_domain(
            rec.length, rec.modification, d.get("mw"), method=rec.method
        )
        rec.confidence, rec.ci_width = confidence_from_ci(rec.ci_low, rec.ci_high)
        # 小分子：置信度来自 B3clf benchmark（~70-73%），而非 CI 宽度
        if rec.domain == "small-molecule" and rec.confidence == "unknown":
            rec.confidence = "benchmarked"
        if not rec.call and rec.p_bbb is not None:
            rec.call = "BBB+" if rec.p_bbb >= 50 else "BBB-"
        return rec

    def to_compact(self) -> dict[str, Any]:
        return {
            "name": self.name,
            "sequence": self.sequence,
            "length": self.length,
            "modification": self.modification or "none",
            "P(BBB+)%": self.p_bbb,
            "call": self.call,
            "bootstrap_90CI": (
                [self.ci_low, self.ci_high] if self.ci_low is not None else None
            ),
            "confidence": self.confidence,
            "applicability_domain": self.domain,
            "prediction_method": self.method,
            "domain_reason": self.domain_reason,
            "known_central_route_literature": self.known_route or None,
            "receptor": self.receptor or None,
        }


# ─────────────────────────────────────────────────────────────────────────────
# LLM 叙述（Opus 优先，失败回退模板）
# ─────────────────────────────────────────────────────────────────────────────

_SYS_BBB_ANALYST = (
    "You are a senior computational scientist specializing in peptide blood-brain "
    "barrier (BBB) permeability prediction. You interpret sequence-based model "
    "outputs (ESM-2 propensity scores with bootstrap intervals) honestly and "
    "precisely. Principles you NEVER violate:\n"
    "1. The model predicts INTRINSIC, sequence-based penetration PROPENSITY - NOT "
    "in-vivo brain exposure, PK-integrated brain levels, or transport route.\n"
    "2. Separate model outputs from literature biological context; never present "
    "literature route claims as model outputs.\n"
    "3. For out-of-domain inputs (length beyond ~20-30 aa, or lipidation/PEGylation/"
    "cyclization/D-amino acids that ESM-2 cannot encode), the score is a "
    "backbone-level UPPER BOUND, not a confident call.\n"
    "4. Write qualitative prose only. Do NOT include any numbers, percentages, or "
    "confidence intervals in your prose - those are inserted programmatically.\n"
    "5. If applicability_domain == 'small-molecule', the input is a NON-PEPTIDE small "
    "molecule predicted by the dedicated, BENCHMARKED small-molecule tool (B3clf / "
    "CNS-MPO), NOT the ESM sequence model. Treat this as a VALID prediction from the "
    "right tool (cite that it is a benchmarked small-molecule classifier) - do NOT call "
    "it out-of-domain or untrustworthy. The ESM sequence model simply does not apply here.\n"
    "Write concise, professional English for a senior R&D audience."
)

# ── 结构化输出 schema（LLM 只填 prose 字段，所有数字由代码注入）──
_PEPTIDE_JSON_SPEC = (
    'Return a JSON object with EXACTLY these string fields:\n'
    '  "reading"    - 1-2 sentences: what the propensity score means for this peptide '
    '(qualitative: high/moderate/low propensity), no numbers.\n'
    '  "domain"     - 1 sentence: the applicability-domain / confidence implication '
    '(in-domain trustworthy vs edge/out-of-domain upper-bound-only).\n'
    '  "caveat"     - 1 sentence: the standing caveat (propensity != in-vivo brain '
    'exposure or route).\n'
    '  "literature" - optional 1 sentence of known central-route context if provided '
    'in the data, else empty string "".\n'
    'No other keys. No numbers/percentages in any field.'
)

# ── few-shot 黄金范例（锚定文风、深度、诚实口径）──
_FEWSHOT_PEPTIDE = (
    "Example A (in-domain, native short peptide):\n"
    '{"reading": "The frozen ESM-2 head assigns this native short peptide a high '
    'intrinsic BBB-penetration propensity, driven by its backbone sequence rather '
    'than any single descriptor.", "domain": "The peptide sits within the training '
    'length distribution, so the call is interpretable and comparatively '
    'trustworthy.", "caveat": "This reflects intrinsic sequence propensity only, not '
    'measured in-vivo brain exposure or the transport route the molecule actually '
    'uses.", "literature": "The endogenous form is reported to reach central targets '
    'via saturable transport."}\n\n'
    "Example B (out-of-domain, lipidated analog):\n"
    '{"reading": "The sequence head returns only a backbone-level propensity because '
    'the defining fatty-acid/PEG modification is invisible to ESM-2.", "domain": '
    '"This input is out-of-domain: the modification cannot be encoded, so the score '
    'is an upper bound, not a confident prediction.", "caveat": "Protraction lowers '
    'free fraction and brain entry in practice, so intrinsic propensity must not be '
    'read as brain exposure.", "literature": ""}'
)

# 校验用：prose 中禁止出现的"数字幻觉"模式（百分数 / 置信区间）
_PCT_RE = re.compile(r"\d+(?:\.\d+)?\s*%")
_CI_RE = re.compile(r"\[\s*\d")


def _chat_json(client, system: str, user: str, *, max_tokens: int = 500) -> dict | None:
    """请求 JSON 模式并解析；失败返回 None。"""
    if client is None or not getattr(client, "enabled", False):
        return None
    raw = client.chat(system, user, max_tokens=max_tokens,
                      response_format={"type": "json_object"})
    if not raw:
        return None
    try:
        obj = json.loads(raw)
        return obj if isinstance(obj, dict) else None
    except (ValueError, TypeError):
        # 容错：截取首个 {...} 块再解析
        m = re.search(r"\{.*\}", raw, re.DOTALL)
        if m:
            try:
                return json.loads(m.group(0))
            except (ValueError, TypeError):
                return None
        return None


def _validate_peptide_fields(fields: dict, rec: "PeptideRecord") -> tuple[bool, str]:
    """确定性校验 LLM 返回的 prose 字段。返回 (ok, reason)。"""
    if not isinstance(fields, dict):
        return False, "not a dict"
    for k in ("reading", "domain", "caveat"):
        v = fields.get(k)
        if not isinstance(v, str) or not v.strip():
            return False, f"missing/empty field: {k}"
        if len(v) > 600:
            return False, f"field too long: {k}"
    prose = " ".join(str(fields.get(k, "")) for k in ("reading", "domain", "caveat", "literature"))
    # 数字归代码：prose 不得含百分数 / 置信区间（防数字幻觉）
    if _PCT_RE.search(prose) or _CI_RE.search(prose):
        return False, "prose contains numeric score/CI (must be code-owned)"
    # 必含 propensity 概念（小分子除外：B3clf 是分类概率，非 ESM propensity）
    if rec.domain != "small-molecule" and "propensity" not in prose.lower():
        return False, "missing 'propensity'"
    # OOD 必须点明"上界/域外/外推"
    if rec.domain.startswith("out-of-domain") or rec.domain == "edge-extrapolation":
        low = prose.lower()
        if not any(t in low for t in ("upper bound", "out-of-domain", "out of domain",
                                      "extrapolat", "backbone")):
            return False, "OOD entry missing upper-bound/OOD note"
    return True, "ok"


def _assemble_peptide_narrative(fields: dict, rec: "PeptideRecord") -> str:
    """用确定性事实 + 校验过的 LLM prose 拼装最终叙述（数字全部由代码注入）。"""
    if rec.domain == "small-molecule":
        if rec.p_bbb is not None:
            lead = f"[MODEL] {rec.name}: B3clf small-molecule P(BBB+) = {rec.p_bbb:.1f}% ({rec.call})."
        else:
            lead = f"[MODEL] {rec.name}: small molecule — use the B3clf small-molecule tool."
    elif rec.p_bbb is not None:
        ci = (f", 90% bootstrap CI [{rec.ci_low:.1f}, {rec.ci_high:.1f}]"
              if rec.ci_low is not None else "")
        lead = f"[MODEL] {rec.name}: ESM-2 P(BBB+) = {rec.p_bbb:.1f}% ({rec.call}){ci}."
    else:
        lead = (f"[MODEL] {rec.name}: not scored by the sequence model "
                f"(out-of-domain - {rec.domain}).")
    parts = [lead, fields["reading"].strip(), fields["domain"].strip(),
             "[CAVEAT] " + fields["caveat"].strip()]
    lit = str(fields.get("literature", "")).strip()
    if lit:
        parts.append("[LITERATURE] " + lit)
    return " ".join(p for p in parts if p)


def generate_peptide_narrative(client, rec: PeptideRecord, *, max_retries: int = 2) -> str:
    """单条肽解读：结构化 JSON + 数字校验 + 重试，失败回退确定性模板。"""
    if client is None or not getattr(client, "enabled", False):
        return _fallback_peptide_narrative(rec)

    base_user = (
        _PEPTIDE_JSON_SPEC + "\n\n" + _FEWSHOT_PEPTIDE
        + "\n\nNow interpret THIS peptide (write prose only, no numbers):\n"
        + json.dumps(rec.to_compact(), ensure_ascii=False, indent=2)
    )
    last_reason = ""
    for attempt in range(max_retries + 1):
        user = base_user
        if last_reason:
            user += (f"\n\nYour previous attempt was rejected: {last_reason}. "
                     "Fix this and return corrected JSON (prose only, no numbers/percentages).")
        fields = _chat_json(client, _SYS_BBB_ANALYST, user, max_tokens=500)
        if fields is not None:
            ok, last_reason = _validate_peptide_fields(fields, rec)
            if ok:
                return _assemble_peptide_narrative(fields, rec)
            logger.info("peptide narrative validation failed (%s), retrying", last_reason)
        else:
            last_reason = "output was not valid JSON"
    logger.warning("peptide narrative fell back to template for %s", rec.name)
    return _fallback_peptide_narrative(rec)


# ── executive summary：同样走结构化 + 校验 + 重试 ──
_EXEC_JSON_SPEC = (
    'Return a JSON object with EXACTLY these string fields:\n'
    '  "overview"    - 1-2 sentences on the overall pattern across the batch.\n'
    '  "trustworthy" - 1 sentence naming which calls are in-domain / trustworthy.\n'
    '  "caution"     - 1 sentence naming which are upper-bound-only '
    '(out-of-domain / modified / extrapolated).\n'
    '  "caveat"      - 1 sentence: scores are intrinsic sequence propensity, not '
    'in-vivo brain exposure.\n'
    'No other keys. You MAY name peptides, but include no numeric scores/CIs.'
)


def _validate_exec_fields(fields: dict) -> tuple[bool, str]:
    if not isinstance(fields, dict):
        return False, "not a dict"
    for k in ("overview", "trustworthy", "caution", "caveat"):
        v = fields.get(k)
        if not isinstance(v, str) or not v.strip():
            return False, f"missing/empty field: {k}"
    prose = " ".join(str(fields.get(k, "")) for k in ("overview", "trustworthy", "caution", "caveat"))
    if _PCT_RE.search(prose) or _CI_RE.search(prose):
        return False, "prose contains numeric score/CI"
    low = prose.lower()
    if "propensity" not in low:
        return False, "missing 'propensity'"
    return True, "ok"


def generate_exec_summary(client, records: list[PeptideRecord], *, max_retries: int = 2) -> str:
    """整批执行摘要：结构化 JSON + 校验 + 重试，失败回退模板。"""
    if client is None or not getattr(client, "enabled", False):
        return _fallback_exec_summary(records)

    rows = [r.to_compact() for r in records]
    base_user = (
        _EXEC_JSON_SPEC
        + "\n\nBatch data (write prose only, no numeric scores):\n"
        + json.dumps(rows, ensure_ascii=False, indent=2)
    )
    last_reason = ""
    for attempt in range(max_retries + 1):
        user = base_user
        if last_reason:
            user += (f"\n\nYour previous attempt was rejected: {last_reason}. "
                     "Fix this and return corrected JSON (prose only, no numbers).")
        fields = _chat_json(client, _SYS_BBB_ANALYST, user, max_tokens=700)
        if fields is not None:
            ok, last_reason = _validate_exec_fields(fields)
            if ok:
                return " ".join(fields[k].strip() for k in
                                ("overview", "trustworthy", "caution", "caveat"))
            logger.info("exec summary validation failed (%s), retrying", last_reason)
        else:
            last_reason = "output was not valid JSON"
    logger.warning("exec summary fell back to template")
    return _fallback_exec_summary(records)


def _fallback_peptide_narrative(rec: PeptideRecord) -> str:
    parts: list[str] = []
    # 小分子：B3clf 有效预测（非 ESM 序列模型，非 OOD）
    if rec.domain == "small-molecule":
        if rec.p_bbb is not None:
            parts.append(
                f"[MODEL] {rec.name}: B3clf small-molecule P(BBB+) = {rec.p_bbb:.1f}% "
                f"({rec.call})."
            )
        parts.append(
            "Applicability: non-peptide small molecule, predicted by the dedicated, "
            "benchmarked B3clf / CNS-MPO tool (the ESM sequence model does not apply here)."
        )
        if rec.known_route:
            parts.append(f"[LITERATURE] {rec.known_route}.")
        parts.append(
            "Caveat: an in-silico BBB classification; confirm with brain PK (Kp,uu) for borderline calls."
        )
        return " ".join(parts)

    if rec.p_bbb is not None:
        parts.append(
            f"[MODEL] {rec.name}: ESM-2 P(BBB+) = {rec.p_bbb:.1f}% ({rec.call})"
            + (
                f", 90% bootstrap CI [{rec.ci_low:.1f}, {rec.ci_high:.1f}]"
                if rec.ci_low is not None
                else ""
            )
            + "."
        )
    dom_txt = {
        "in-domain": "input within training distribution; score is interpretable as sequence propensity.",
        "edge-extrapolation": "input near training maximum; score is edge extrapolation with reduced confidence.",
        "out-of-domain-length": "length beyond training maximum; score is extrapolation, not a confident call.",
        "out-of-domain-modification": "modification cannot be encoded by ESM-2; score is a backbone-level upper bound only.",
        "out-of-domain-modality": "neither a standard peptide nor covered by the small-molecule tool; no validated predictor applies.",
    }.get(rec.domain, "")
    if dom_txt:
        parts.append(f"Applicability: {rec.domain} ({dom_txt})")
    if rec.confidence in {"low", "medium"} and rec.ci_width is not None:
        parts.append(f"Confidence {rec.confidence} (CI width {rec.ci_width:.0f} pts).")
    if rec.known_route:
        parts.append(f"[LITERATURE] Known central route: {rec.known_route}.")
    parts.append(
        "Caveat: this score is intrinsic sequence propensity, not in-vivo brain "
        "exposure or transport route."
    )
    return " ".join(parts)


def _fallback_exec_summary(records: list[PeptideRecord]) -> str:
    n = len(records)
    in_dom = [r.name for r in records if r.domain == "in-domain"]
    sm = [r.name for r in records if r.domain == "small-molecule"]
    ood = [r.name for r in records if r.domain.startswith("out-of-domain")]
    pos = [r.name for r in records if r.call == "BBB+"]
    lines = [
        f"This batch covers {n} entries.",
        f"Peptide in-domain (trustworthy): {len(in_dom)} - {', '.join(in_dom) if in_dom else 'none'}.",
    ]
    if sm:
        lines.append(
            f"Small molecules via the dedicated B3clf tool (not the sequence model): "
            f"{len(sm)} - {', '.join(sm)}."
        )
    lines += [
        f"Out-of-domain for any validated predictor: {len(ood)} - {', '.join(ood) if ood else 'none'}.",
        f"Predicted BBB+: {len(pos)}.",
        "Peptide scores are intrinsic sequence-based penetration propensity, not "
        "in-vivo brain exposure; small-molecule calls come from the benchmarked B3clf tool.",
    ]
    return " ".join(lines)


# ─────────────────────────────────────────────────────────────────────────────
# 自包含 HTML 渲染（navy 主题 + CSS 条形图 + 诚实口径）
# ─────────────────────────────────────────────────────────────────────────────

_DOMAIN_BADGE = {
    "in-domain": ("in-domain", "#0a8f5b", "#e3f3ea"),
    "edge-extrapolation": ("edge", "#9a6a00", "#f5edda"),
    "out-of-domain-length": ("OOD · length", "#b3261e", "#f7e4e2"),
    "out-of-domain-modification": ("OOD · modification", "#b3261e", "#f7e4e2"),
    "out-of-domain-modality": ("OOD · non-peptide", "#7a1f6e", "#f3e2f0"),
    "small-molecule": ("small molecule · B3clf", "#0a6e8f", "#e2f0f5"),
}
_CONF_BADGE = {"high": "#0a8f5b", "medium": "#9a6a00", "low": "#b3261e", "unknown": "#5b6470", "benchmarked": "#0a6e8f"}


def _adaptive_reading_note(records: list[PeptideRecord]) -> str:
    """Build modality-aware reading note with explicit uncertainty capability limits."""
    has_peptide = any(r.domain != "small-molecule" for r in records)
    has_sm = any(r.domain == "small-molecule" for r in records)

    common = (
        "Interpret [LITERATURE] route context as external biology, not a model output. "
        "Propensity/classification is not equivalent to in-vivo brain exposure."
    )

    if has_peptide and not has_sm:
        return (
            "Peptide branch: scores are intrinsic, sequence-based BBB penetration propensity "
            "from a frozen ESM-2 model. Uncertainty is estimated by bootstrap 90% CI, which "
            "captures sampling variance but does not capture domain shift or PK effects. "
            "ESM-2 reads only the 20 standard amino acids, so lipidation / PEGylation / "
            "cyclization / D-amino acids are not encoded; such entries are backbone-level upper bounds. "
            + common
        )

    if has_sm and not has_peptide:
        return (
            "Small-molecule branch: calls come from the dedicated B3clf / CNS-MPO path, not ESM-2. "
            "Current uncertainty is benchmark-level only (~70-73% historical performance); "
            "there is no calibrated per-compound confidence interval in this branch yet. "
            "Treat borderline probabilities as uncertain and prioritize brain PK confirmation (for example Kp,uu). "
            + common
        )

    return (
        "Mixed batch: peptide entries use ESM-2 sequence propensity with bootstrap 90% CI as the explicit "
        "uncertainty signal (sampling variance only), while small-molecule entries use B3clf / CNS-MPO with "
        "benchmark-level confidence only and no calibrated per-compound CI yet. "
        "For modified peptides beyond ESM encoding, read scores as backbone upper bounds. "
        "For borderline small-molecule probabilities, prioritize brain PK confirmation. "
        + common
    )


def _esc(s: Any) -> str:
    return _html.escape(str(s if s is not None else ""))


def _bar(p: float | None) -> str:
    if p is None:
        return '<span class="muted">—</span>'
    pos = p >= 50
    cls = "bg-pos" if pos else "bg-neg"
    label_pos = "right:6px" if pos else f"left:calc({min(p,100):.1f}% + 6px)"
    return (
        f'<div class="bar"><span class="{cls}" style="width:{min(max(p,0),100):.1f}%"></span>'
        f'<span class="lab" style="{label_pos}">{p:.1f}%</span></div>'
    )


def _display_len(rec: PeptideRecord) -> str:
    if rec.length is not None:
        return str(rec.length)
    if rec.domain == "small-molecule":
        return "SM"
    return "-"


def _adaptive_ci_header(records: list[PeptideRecord]) -> str:
    """Generate modality-aware CI column header."""
    has_peptide = any(r.domain != "small-molecule" for r in records)
    has_sm = any(r.domain == "small-molecule" for r in records)
    if has_peptide and has_sm:
        return "Uncertainty (90% CI / Vote-CI90)"
    if has_sm:
        return "Vote-CI90"
    return "90% CI"


def _display_ci(rec: PeptideRecord) -> str:
    if rec.ci_low is not None and rec.ci_high is not None:
        return f"[{rec.ci_low:.1f}, {rec.ci_high:.1f}]"
    if rec.domain == "small-molecule":
        if rec.committee_pos is not None and rec.committee_n:
            lo, hi = _wilson_interval(rec.committee_pos, rec.committee_n)
            return (
                f"Vote-CI90 [{lo*100:.1f}, {hi*100:.1f}] "
                f"({rec.committee_pos}/{rec.committee_n})"
            )
        return "N/A (benchmarked)"
    return "N/A"


def render_html(
    records: list[PeptideRecord],
    exec_summary: str,
    *,
    title: str = "BBB Penetration Prediction — Auto Report",
    subtitle: str = "ESM-2 sequence propensity · Opus-generated interpretation",
    signature: str = "NNRCC agenter group",
    llm_status: str = "",
) -> str:
    """渲染自包含 HTML 字符串。"""
    date = _dt.date.today().isoformat()
    reading_note = _adaptive_reading_note(records)

    rows_html = []
    for r in records:
        badge_txt, badge_fg, badge_bg = _DOMAIN_BADGE.get(
            r.domain, (r.domain or "—", "#5b6470", "#eef1f6")
        )
        ci_txt = _display_ci(r)
        len_txt = _display_len(r)
        conf_fg = _CONF_BADGE.get(r.confidence, "#5b6470")
        narrative_html = _esc(r.narrative).replace("\n", "<br>")
        # 高亮来源标签
        narrative_html = (
            narrative_html.replace("[MODEL]", '<b style="color:#0058a3">[MODEL]</b>')
            .replace("[LITERATURE]", '<b style="color:#9a6a00">[LITERATURE]</b>')
        )
        rows_html.append(
            f"""
      <tr>
        <td><b>{_esc(r.name)}</b><div class="seq">{_esc(r.sequence)}</div></td>
                <td class="num">{_esc(len_txt)}</td>
        <td style="min-width:140px">{_bar(r.p_bbb)}</td>
        <td class="num muted">{_esc(ci_txt)}</td>
        <td><span class="badge" style="color:{conf_fg};border-color:{conf_fg}">{_esc(r.confidence)}</span></td>
        <td><span class="badge" style="color:{badge_fg};background:{badge_bg};border-color:{badge_fg}">{_esc(badge_txt)}</span></td>
      </tr>
      <tr class="narr"><td colspan="6">{narrative_html}</td></tr>"""
        )

    status_note = f' · {_esc(llm_status)}' if llm_status else ""

    return f"""<!DOCTYPE html>
<html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{_esc(title)}</title>
<style>
  :root{{--blue:#001965;--accent:#0058a3;--light:#e8eef7;--grey:#5b6470;--green:#0a8f5b;--red:#b3261e;--line:#dde3ec;}}
  *{{box-sizing:border-box}}
  body{{font-family:"Segoe UI",-apple-system,"Noto Sans",Helvetica,Arial,sans-serif;color:#1a1f2e;background:#eef1f6;margin:0;line-height:1.5;font-size:15px;}}
  .wrap{{max-width:1000px;margin:0 auto;padding:0 18px 60px;}}
  header.top{{background:linear-gradient(135deg,#001965 0%,#0058a3 100%);color:#fff;padding:32px 40px 28px;border-radius:0 0 14px 14px;}}
  header.top h1{{margin:0 0 6px;font-size:26px;font-weight:700;}}
  header.top .sub{{font-size:15px;color:#cfe0f4;}}
  header.top .meta{{margin-top:14px;font-size:12.5px;color:#bcd0ec;}}
  .lead{{background:#fff;border-left:6px solid var(--accent);border-radius:8px;padding:16px 20px;margin:22px 0;box-shadow:0 1px 4px rgba(0,25,101,.07);}}
  .lead h2{{margin:0 0 8px;font-size:15px;color:var(--blue);text-transform:uppercase;letter-spacing:.5px;}}
  h2.sec{{color:var(--blue);font-size:19px;margin:30px 0 8px;padding-bottom:6px;border-bottom:2px solid var(--light);}}
  table{{border-collapse:collapse;width:100%;font-size:13.5px;background:#fff;border-radius:9px;overflow:hidden;box-shadow:0 1px 4px rgba(0,25,101,.07);}}
  th{{background:var(--blue);color:#fff;text-align:left;padding:9px 11px;font-weight:600;}}
  td{{padding:8px 11px;border-bottom:1px solid var(--line);vertical-align:middle;}}
  td.num{{text-align:right;font-variant-numeric:tabular-nums;}}
  .seq{{font-family:ui-monospace,Menlo,Consolas,monospace;font-size:11px;color:var(--grey);margin-top:2px;word-break:break-all;}}
  .muted{{color:var(--grey);}}
  tr.narr td{{background:#f6f9fc;font-size:13px;color:#33384a;padding:10px 14px 14px;border-bottom:2px solid var(--line);}}
  .bar{{position:relative;height:18px;background:#eef1f6;border-radius:4px;min-width:120px;overflow:hidden;}}
  .bar>span{{position:absolute;left:0;top:0;bottom:0;border-radius:4px;}}
  .bar .lab{{position:absolute;top:0;bottom:0;display:flex;align-items:center;font-size:11.5px;font-weight:700;color:#1a1f2e;}}
  .bg-pos{{background:linear-gradient(90deg,#5bbf8f,#0a8f5b);}}
  .bg-neg{{background:linear-gradient(90deg,#e39b95,#b3261e);}}
  .badge{{display:inline-block;font-size:11px;font-weight:700;border:1px solid;border-radius:4px;padding:2px 7px;white-space:nowrap;}}
  .caveat{{background:#fff8e6;border-left:5px solid #9a6a00;border-radius:8px;padding:13px 18px;margin:20px 0;font-size:13.5px;}}
  footer{{margin-top:28px;font-size:12.5px;color:var(--grey);text-align:center;}}
  @media print{{body{{background:#fff}}.lead,table{{box-shadow:none;border:1px solid var(--line)}}header.top{{border-radius:0}}}}
</style></head>
<body>
<header class="top">
  <h1>{_esc(title)}</h1>
  <div class="sub">{_esc(subtitle)}</div>
  <div class="meta">{_esc(signature)} · {date}{status_note}</div>
</header>
<div class="wrap">
  <div class="lead">
    <h2>Executive summary</h2>
    <div>{_esc(exec_summary).replace(chr(10), '<br>')}</div>
  </div>

  <h2 class="sec">Per-peptide predictions &amp; interpretation</h2>
  <table>
        <tr><th>Entry</th><th>Len</th><th>P(BBB+)%</th><th>{_adaptive_ci_header(records)}</th><th>Conf.</th><th>Applicability</th></tr>
    {''.join(rows_html)}
  </table>

  <div class="caveat">
        <b>Reading note.</b> {_esc(reading_note)}
  </div>

  <footer>Auto-generated by bbbkit.peptide.report · interpretation by Claude Opus 4.7 (graceful-degrades to templates).</footer>
</div>
</body></html>"""


# ─────────────────────────────────────────────────────────────────────────────
# PPTX 渲染（python-pptx，可选依赖；缺失时给出清晰报错）
# ─────────────────────────────────────────────────────────────────────────────

_PPTX_NAVY = (0x00, 0x19, 0x65)
_PPTX_ACCENT = (0x00, 0x58, 0xA3)
_PPTX_GREEN = (0x0A, 0x8F, 0x5B)
_PPTX_RED = (0xB3, 0x26, 0x1E)
_PPTX_AMBER = (0x9A, 0x6A, 0x00)
_PPTX_GREY = (0x5B, 0x64, 0x70)
_PPTX_WHITE = (0xFF, 0xFF, 0xFF)
_PPTX_LIGHT = (0xE8, 0xEE, 0xF7)
_PPTX_INK = (0x1A, 0x1F, 0x2E)

_DOMAIN_SHORT = {
    "in-domain": ("in-domain", _PPTX_GREEN),
    "edge-extrapolation": ("edge", _PPTX_AMBER),
    "out-of-domain-length": ("OOD - length", _PPTX_RED),
    "out-of-domain-modification": ("OOD - modification", _PPTX_RED),
    "out-of-domain-modality": ("OOD - non-peptide", (0x7A, 0x1F, 0x6E)),
    "small-molecule": ("small molecule - B3clf", (0x0A, 0x6E, 0x8F)),
}


def _scrub_pptx_cjk_fonts(path: str) -> None:
    """清除 python-pptx 默认主题 theme1.xml 里的东亚字体名（如宋体/新細明體），
    使导出的 PPTX 字节层面不含 CJK 字符。报告内容本身为纯西文，故无副作用。"""
    import os
    import re as _re
    import tempfile
    import zipfile

    with zipfile.ZipFile(path) as src:
        items = [(i, src.read(i.filename)) for i in src.infolist()]
    # 临时文件必须与目标同目录，否则 os.replace 跨文件系统会失败（Errno 18）
    dirpath = os.path.dirname(os.path.abspath(path)) or "."
    fd, tmp = tempfile.mkstemp(suffix=".pptx", dir=dirpath)
    os.close(fd)
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for info, data in items:
            if info.filename == "ppt/theme/theme1.xml":
                txt = data.decode("utf-8")
                txt = _re.sub(
                    r'(typeface=")[^"]*[\u4e00-\u9fff][^"]*(")', r"\1\2", txt)
                data = txt.encode("utf-8")
            out.writestr(info, data)
    os.replace(tmp, path)


def render_pptx(
    records: list["PeptideRecord"],
    exec_summary: str,
    out_path,
    *,
    title: str = "BBB Penetration Prediction - Auto Report",
    subtitle: str = "ESM-2 sequence propensity - Opus-generated interpretation",
    signature: str = "NNRCC agenter group",
    llm_status: str = "",
) -> str:
    """渲染 PPTX 幻灯片到 out_path（标题 / 摘要 / 预测表 / 阅读须知）。

    需要 python-pptx（``pip install python-pptx``）；缺失时抛 ImportError。
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Emu, Pt
    except ImportError as exc:  # pragma: no cover - 依赖缺失路径
        raise ImportError(
            "PPTX 导出需要 python-pptx：pip install python-pptx"
        ) from exc

    SLIDE_W, SLIDE_H = 12192000, 6858000
    date = _dt.date.today().isoformat()
    reading_note = _adaptive_reading_note(records)

    def rgb(t):
        return RGBColor(*t)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    blank = prs.slide_layouts[6]

    def add_text(slide, l, t, w, h, text, *, size=18, bold=False,
                 color=_PPTX_NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)
        return box

    def add_rect(slide, l, t, w, h, color):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(color)
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    # ── Slide 1: 标题页 ──
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, _PPTX_NAVY)
    add_rect(s, 0, 4950000, SLIDE_W, 18000, _PPTX_ACCENT)
    add_text(s, 900000, 2100000, SLIDE_W - 1800000, 1300000, title,
             size=38, bold=True, color=_PPTX_WHITE)
    add_text(s, 900000, 3450000, SLIDE_W - 1800000, 800000, subtitle,
             size=19, color=_PPTX_LIGHT)
    add_text(s, 900000, 5550000, SLIDE_W - 1800000, 600000,
             f"{signature}  -  {date}", size=14, color=_PPTX_LIGHT)

    # ── Slide 2: 执行摘要 ──
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, 950000, _PPTX_NAVY)
    add_text(s, 600000, 250000, SLIDE_W - 1200000, 600000, "Executive summary",
             size=26, bold=True, color=_PPTX_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 600000, 1250000, SLIDE_W - 1200000, 5000000, exec_summary,
             size=16, color=_PPTX_INK)

    # ── Slide 3: 预测汇总表 ──
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, 950000, _PPTX_NAVY)
    add_text(s, 600000, 250000, SLIDE_W - 1200000, 600000,
             "Per-peptide predictions", size=26, bold=True, color=_PPTX_WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    ci_header = _adaptive_ci_header(records)
    headers = ["Entry", "Len", "P(BBB+)%", ci_header, "Conf.", "Applicability"]
    widths = [3000000, 900000, 1700000, 2200000, 1400000, 2600000]
    nrows = len(records) + 1
    tbl = s.shapes.add_table(
        nrows, len(headers), Emu(600000), Emu(1150000),
        Emu(sum(widths)), Emu(min(5200000, 360000 * nrows))).table
    for j, w in enumerate(widths):
        tbl.columns[j].width = Emu(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = rgb(_PPTX_NAVY)
        pr = c.text_frame.paragraphs[0].runs[0]
        pr.font.size = Pt(11)
        pr.font.bold = True
        pr.font.color.rgb = rgb(_PPTX_WHITE)
    for i, rec in enumerate(records, start=1):
        dom_txt, dom_col = _DOMAIN_SHORT.get(rec.domain, (rec.domain or "-", _PPTX_GREY))
        ci = _display_ci(rec)
        len_txt = _display_len(rec)
        p_txt = f"{rec.p_bbb:.1f}" if rec.p_bbb is not None else "-"
        vals = [rec.name, len_txt, p_txt, ci,
                rec.confidence or "-", dom_txt]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.text = v
            pr = c.text_frame.paragraphs[0].runs[0]
            pr.font.size = Pt(10)
            if j == 5:
                pr.font.color.rgb = rgb(dom_col)
                pr.font.bold = True
            elif j == 2 and rec.p_bbb is not None:
                pr.font.color.rgb = rgb(_PPTX_GREEN if rec.p_bbb >= 50 else _PPTX_RED)
                pr.font.bold = True
            else:
                pr.font.color.rgb = rgb(_PPTX_INK)

    # ── Slide 4: 阅读须知（诚实口径）──
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, 950000, _PPTX_AMBER)
    add_text(s, 600000, 250000, SLIDE_W - 1200000, 600000, "Reading note",
             size=26, bold=True, color=_PPTX_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    note = reading_note
    add_text(s, 600000, 1300000, SLIDE_W - 1200000, 4400000, note, size=16, color=_PPTX_INK)
    if llm_status:
        add_text(s, 600000, 6300000, SLIDE_W - 1200000, 350000,
                 f"Interpretation: {llm_status}", size=11, color=_PPTX_GREY)

    out_path = str(out_path)
    prs.save(out_path)
    _scrub_pptx_cjk_fonts(out_path)
    return out_path


# ─────────────────────────────────────────────────────────────────────────────
# 端到端编排
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class BBBReportBundle:
    outdir: str
    html_path: str = ""
    pptx_path: str = ""
    matrix_csv: str = ""
    records: list[dict[str, Any]] = field(default_factory=list)
    llm_status: str = ""

    def to_dict(self) -> dict:
        return {
            "outdir": self.outdir,
            "html_path": self.html_path,
            "pptx_path": self.pptx_path,
            "matrix_csv": self.matrix_csv,
            "n_peptides": len(self.records),
            "llm_status": self.llm_status,
        }


def _get_client(use_llm: bool):
    if not use_llm:
        return None
    try:
        from ..report.llm import LLMClient

        return LLMClient()
    except Exception as e:  # noqa: BLE001
        logger.warning("LLMClient init failed, falling back to templates: %s", e)
        return None


def _write_matrix_csv(records: list[PeptideRecord], path: Path) -> None:
    import csv

    cols = [
        "name", "sequence", "length", "modification", "p_bbb", "call",
        "ci_low", "ci_high", "confidence", "domain", "domain_reason", "known_route",
    ]
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=cols)
        w.writeheader()
        for r in records:
            w.writerow({c: getattr(r, c, "") for c in cols})


def build_bbb_report(
    predictions: list[dict[str, Any]],
    outdir: str,
    *,
    title: str = "BBB Penetration Prediction — Auto Report",
    subtitle: str = "ESM-2 sequence propensity · Opus-generated interpretation",
    signature: str = "NNRCC agenter group",
    use_llm: bool = True,
    pptx: bool = False,
) -> BBBReportBundle:
    """端到端：预测记录 → 适用域标记 → LLM 叙述 → 自包含 HTML + 矩阵 CSV。

    Parameters
    ----------
    predictions : list[dict]
        每条含 name + (sequence 或 length) + p_bbb/ESM2_P(BBB+)%；可选
        modification / ci / known_route / receptor / call。
    outdir : str
        输出目录（自动创建）。
    use_llm : bool
        True 时调用 Opus 4.7 生成叙述；无配置/失败自动回退模板。
    """
    out = Path(outdir)
    out.mkdir(parents=True, exist_ok=True)

    records = [PeptideRecord.from_dict(d) for d in predictions]

    client = _get_client(use_llm)
    llm_status = getattr(client, "status", "LLM disabled — templates") if client else "LLM disabled — templates"
    logger.info("LLM: %s", llm_status)

    for rec in records:
        rec.narrative = generate_peptide_narrative(client, rec)
    exec_summary = generate_exec_summary(client, records)

    html = render_html(
        records, exec_summary, title=title, subtitle=subtitle,
        signature=signature, llm_status=llm_status,
    )
    html_path = out / "bbb_auto_report.html"
    html_path.write_text(html, encoding="utf-8")

    matrix_csv = out / "bbb_predictions_matrix.csv"
    _write_matrix_csv(records, matrix_csv)

    pptx_path = ""
    if pptx:
        try:
            pptx_path = render_pptx(
                records, exec_summary, out / "bbb_auto_report.pptx",
                title=title, subtitle=subtitle, signature=signature,
                llm_status=llm_status,
            )
        except ImportError as e:  # noqa: BLE001
            logger.warning("PPTX export skipped: %s", e)

    return BBBReportBundle(
        outdir=str(out),
        html_path=str(html_path),
        pptx_path=str(pptx_path),
        matrix_csv=str(matrix_csv),
        records=[r.to_compact() for r in records],
        llm_status=llm_status,
    )
