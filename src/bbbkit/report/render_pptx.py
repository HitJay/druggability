"""
bbbkit.report.render_pptx — python-pptx 幻灯片渲染

生成幻灯片:
  1. 标题页
  2. 执行摘要 (LLM 叙述)
  3. 二维定位散点 (validation × tractability，用 pptx shapes 绘制)
  4. 对比矩阵表
  5. 每靶一页 (关键指标 + LLM 叙述)
"""

from __future__ import annotations

import datetime as _dt
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# 16:9 画布 (EMU): 12192000 x 6858000
_SLIDE_W = 12192000
_SLIDE_H = 6858000

_ACCENT = RGBColor(0x15, 0x65, 0xC0)
_DARK = RGBColor(0x1A, 0x1A, 0x1A)
_MUTED = RGBColor(0x66, 0x66, 0x66)
_MOD_COLORS = {
    "small_molecule": RGBColor(0x15, 0x65, 0xC0),
    "antibody": RGBColor(0xC6, 0x28, 0x28),
    "protac": RGBColor(0x6A, 0x1B, 0x9A),
}
_QUAD = {
    "priority": RGBColor(0xE8, 0xF5, 0xE9),
    "hard": RGBColor(0xFF, 0xF8, 0xE1),
    "tract": RGBColor(0xE3, 0xF2, 0xFD),
    "watch": RGBColor(0xFA, 0xFA, 0xFA),
}


def _mod_label(m: str | None) -> str:
    return {"small_molecule": "小分子", "antibody": "抗体", "protac": "PROTAC/降解剂",
            "unknown": "未知", "(offline)": "—"}.get(m or "", m or "—")


def _add_textbox(slide, l, t, w, h, text, *, size=18, bold=False, color=_DARK,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_multiline(slide, l, t, w, h, text, *, size=14, color=_DARK):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _verdict_color(rec: str | None) -> RGBColor:
    rec = rec or ""
    if rec.startswith("Priority"):
        return RGBColor(0x2E, 0x7D, 0x32)
    if rec.startswith("Hard"):
        return RGBColor(0xF9, 0xA8, 0x25)
    if rec.startswith("Tractable"):
        return RGBColor(0x15, 0x65, 0xC0)
    return _MUTED


def _slide_title(prs, title: str, subtitle: str = "") -> None:
    s = _blank(prs)
    _add_textbox(s, 700000, 2400000, 10800000, 900000, title, size=40, bold=True, color=_ACCENT)
    if subtitle:
        _add_textbox(s, 700000, 3400000, 10800000, 600000, subtitle, size=18, color=_MUTED)


def _slide_summary(prs, text: str) -> None:
    s = _blank(prs)
    _add_textbox(s, 600000, 350000, 11000000, 700000, "执行摘要", size=28, bold=True, color=_ACCENT)
    _add_multiline(s, 600000, 1250000, 11000000, 5200000, text, size=16)


def _slide_portfolio(prs, records: list[dict[str, Any]]) -> None:
    s = _blank(prs)
    _add_textbox(s, 600000, 300000, 11000000, 700000, "二维定位 (遗传学验证 × 可药性)", size=28, bold=True, color=_ACCENT)
    # plot area
    px, py = 1500000, 1200000
    pw, ph = 8200000, 4600000
    xt = px + pw * 0.6  # tractability threshold 0.6
    yt = py + ph * (1 - 0.55)  # genetics threshold 0.55 (y inverted)
    from pptx.enum.shapes import MSO_SHAPE

    def rect(l, t, w, h, color):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(l)), Emu(int(t)), Emu(int(w)), Emu(int(h)))
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        sp.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); sp.line.width = Pt(0.5)
        sp.shadow.inherit = False
        return sp

    rect(xt, py, px + pw - xt, yt - py, _QUAD["priority"])
    rect(px, py, xt - px, yt - py, _QUAD["hard"])
    rect(xt, yt, px + pw - xt, py + ph - yt, _QUAD["tract"])
    rect(px, yt, xt - px, py + ph - yt, _QUAD["watch"])
    # quadrant labels
    _add_textbox(s, int(xt) + 60000, int(py) + 40000, 2500000, 350000, "Priority", size=12, color=RGBColor(0x2E, 0x7D, 0x32))
    _add_textbox(s, int(px) + 60000, int(py) + 40000, 2500000, 350000, "Hard but worth it", size=12, color=RGBColor(0xF9, 0xA8, 0x25))
    _add_textbox(s, int(xt) + 60000, int(yt) - 360000, 2500000, 350000, "Tractable but verify", size=12, color=_ACCENT)
    # axis labels
    _add_textbox(s, px, py + ph + 80000, pw, 350000, "Tractability (可药性，最优模态分) →", size=12, color=_MUTED, align=PP_ALIGN.CENTER)
    _add_textbox(s, 300000, py, 900000, ph, "Genetics 验证 →", size=12, color=_MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # data points
    from pptx.enum.shapes import MSO_SHAPE as _SH
    r_marker = 150000
    for rec in records:
        gx, gy = rec.get("tractability_best"), rec.get("genetics_score")
        if gx is None or gy is None:
            continue
        cx = px + pw * max(0.0, min(1.0, float(gx)))
        cy = py + ph * (1 - max(0.0, min(1.0, float(gy))))
        color = _MOD_COLORS.get(rec.get("best_modality", ""), _MUTED)
        dot = s.shapes.add_shape(_SH.OVAL, Emu(int(cx - r_marker)), Emu(int(cy - r_marker)),
                                 Emu(r_marker * 2), Emu(r_marker * 2))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); dot.line.width = Pt(1)
        dot.shadow.inherit = False
        _add_textbox(s, int(cx) + 180000, int(cy) - 180000, 2000000, 350000,
                     str(rec.get("gene_name", "")), size=12, bold=True)
    # legend
    lx = px + pw + 250000
    for i, (k, lbl) in enumerate([("small_molecule", "小分子"), ("antibody", "抗体"), ("protac", "PROTAC")]):
        ly = py + i * 450000
        dot = s.shapes.add_shape(_SH.OVAL, Emu(int(lx)), Emu(int(ly)), Emu(220000), Emu(220000))
        dot.fill.solid(); dot.fill.fore_color.rgb = _MOD_COLORS[k]
        dot.line.fill.background(); dot.shadow.inherit = False
        _add_textbox(s, int(lx) + 320000, int(ly) - 40000, 1500000, 350000, lbl, size=12)


def _slide_matrix(prs, records: list[dict[str, Any]]) -> None:
    s = _blank(prs)
    _add_textbox(s, 500000, 250000, 11000000, 600000, "对比矩阵", size=28, bold=True, color=_ACCENT)
    cols = ["靶点", "GWAS", "genetics", "tract", "模态", "ligand", "overall", "结论"]
    rows = len(records) + 1
    tbl_shape = s.shapes.add_table(rows, len(cols), Emu(500000), Emu(1050000), Emu(11200000), Emu(700000 + rows * 360000))
    table = tbl_shape.table
    for j, c in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = c
        cell.fill.solid(); cell.fill.fore_color.rgb = _ACCENT
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(12); run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, r in enumerate(records, 1):
        vals = [
            str(r.get("gene_name", "")),
            str(r.get("gwas_trait", "")),
            _fmt(r.get("genetics_score")),
            _fmt(r.get("tractability_best")),
            _mod_label(r.get("best_modality")),
            _fmt(r.get("ligandability_score")),
            _fmt(r.get("overall_score")),
            (r.get("recommendation") or "").split("—")[0].strip(),
        ]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = v or "—"  # 空串不会生成 run → 用 "—" 占位
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(11)
            if j == 7:
                run.font.color.rgb = _verdict_color(r.get("recommendation"))
                run.font.bold = True


def _slide_target(prs, r: dict[str, Any]) -> None:
    s = _blank(prs)
    gene = r.get("gene_name", "?")
    _add_textbox(s, 500000, 280000, 9000000, 700000, f"{gene}", size=30, bold=True, color=_ACCENT)
    _add_textbox(s, 500000, 950000, 11000000, 400000,
                 f"{r.get('gene_id','')} · GWAS {r.get('gwas_trait','')} · 类别 {r.get('target_class') or '—'} · top 疾病 {r.get('top_disease') or '—'}",
                 size=13, color=_MUTED)
    # verdict badge
    _add_textbox(s, 9700000, 280000, 2000000, 600000, (r.get("recommendation") or "").split("—")[0].strip(),
                 size=16, bold=True, color=_verdict_color(r.get("recommendation")), align=PP_ALIGN.RIGHT)
    # metrics row
    metrics = [
        ("Genetics", _fmt(r.get("genetics_score"))),
        ("Tractability", _fmt(r.get("tractability_best"))),
        ("最优模态", _mod_label(r.get("best_modality"))),
        ("Ligandability", _fmt(r.get("ligandability_score"))),
        ("Overall", _fmt(r.get("overall_score"))),
    ]
    from pptx.enum.shapes import MSO_SHAPE
    mw = 2150000
    for i, (k, v) in enumerate(metrics):
        l = 500000 + i * (mw + 120000)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(1500000), Emu(mw), Emu(1100000))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFB)
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); card.line.width = Pt(0.5)
        card.shadow.inherit = False
        _add_textbox(s, l + 60000, 1580000, mw - 120000, 350000, k, size=11, color=_MUTED)
        _add_textbox(s, l + 60000, 1900000, mw - 120000, 600000, v, size=22, bold=True)
    # narrative
    narr = r.get("_narrative", "") or r.get("modality_note", "")
    _add_textbox(s, 500000, 2900000, 11200000, 600000, "分析", size=16, bold=True, color=_ACCENT)
    _add_multiline(s, 500000, 3450000, 11200000, 2900000, narr, size=15)


def _fmt(v: Any) -> str:
    if v is None:
        return "—"
    if isinstance(v, float):
        return f"{v:.2f}"
    return str(v)


def render_pptx(
    records: list[dict[str, Any]],
    out_path: str,
    *,
    title: str = "靶点深度可药性评估报告",
    exec_summary: str = "",
) -> str:
    """生成 PPTX 并保存到 out_path，返回路径。"""
    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(_SLIDE_H)

    _slide_title(prs, title, f"生成日期 {_dt.date.today().isoformat()} · {len(records)} 个靶点")
    if exec_summary:
        _slide_summary(prs, exec_summary)
    _slide_portfolio(prs, records)
    _slide_matrix(prs, records)
    for r in records:
        _slide_target(prs, r)

    prs.save(out_path)
    return out_path
